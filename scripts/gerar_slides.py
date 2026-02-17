from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_slide_content():
    """
    Retorna uma lista de dicionários com o conteúdo de cada slide
    baseado no Roteiro de Defesa do TCC SINAPSE-BR IA.
    """
    slides_data = [
        {
            "layout": 0,  # Title Slide
            "title": "Rubrica SINAPSE-BR IA",
            "content": [
                "Superando a Dualidade Estrutural na EPT",
                "Integração: Neurociência, Politecnia e Geofilosofia",
                "Autor: Neirivon Elias Cardoso",
                "Orientadora: Profa. Dra. Thays Martins Vital da Silva",
                "IFTM - 2026"
            ],
            "notes": "Bom dia/boa noite à banca. Apresento a Rubrica SINAPSE-BR IA. Antes dos aspectos técnicos, convido a olhar para a questão central: superar a 'fratura educacional' histórica que divide quem pensa de quem faz.",
            "image_desc": "Fundo sutil com logo do IFTM e rede neural."
        },
        {
            "layout": 1,  # Title and Content
            "title": "O Problema: A Fratura e a Curvatura da Vara",
            "content": [
                "Dualidade Estrutural: Pensar vs. Fazer.",
                "Risco: Exclusão Branda (Aluno presente, mas marginalizado).",
                "Metáfora: 'A Curvatura da Vara' (Saviani).",
                "Objetivo: Equilíbrio entre Técnica e Formação Humana."
            ],
            "notes": "O modelo tradicional classifica e exclui. Usamos a metáfora de Saviani: 'A Curvatura da Vara'. Se a educação está curvada para o técnico fragmentado, precisamos curvá-la conscientemente para o humano para que fique reta.",
            "image_desc": "Imagem ilustrativa de uma vara de madeira sendo corrigida de uma curva excessiva para o centro."
        },
        {
            "layout": 1,
            "title": "Fundamentação: O Chão Geográfico e a Luz",
            "content": [
                "Geofilosofia: O Território como condicionante.",
                "Metáfora 1: 'O Chão Geográfico' (Realidade de Sobradinho).",
                "Metáfora 2: 'A Teoria como Luz' (Freire).",
                "Contexto: Triângulo Mineiro e Alto Paranaíba (TMAP)."
            ],
            "notes": "O pensamento não ocorre no vácuo. Usamos a metáfora do 'Chão Geográfico': a avaliação está enraizada no território. E como iluminamos esse chão? Com a 'Teoria como Luz'. A ciência ilumina a prática, e a prática dá concretude à teoria.",
            "image_desc": "Uma planta de raízes profundas em um mapa do Triângulo Mineiro, iluminada por um feixe de luz."
        },
        {
            "layout": 1,
            "title": "A Conexão: A Metáfora da Sinapse",
            "content": [
                "Projeto: SINAPSE-BR IA.",
                "Metáfora Central: 'A Sinapse Pedagógica'.",
                "Triade: Ensino + Aprendizagem + Tecnologia.",
                "Função: Transformar julgamento estático em comunicação dinâmica."
            ],
            "notes": "Por que Sinapse? Biologicamente, é a comunicação entre neurônios. Pedagogicamente, é o 'Encontro Pedagógico'. É onde a avaliação conecta professor, aluno e IA. A IA é o neurotransmissor que facilita a conexão, não o cérebro que decide.",
            "image_desc": "Esquema de uma sinapse neural, rotulada com: Professor, Aluno, Tecnologia (IA)."
        },
        {
            "layout": 1,
            "title": "O Produto: Voxel e Permeabilidade Seletiva",
            "content": [
                "Metáfora 3: 'O Voxel' (Marcador 3D de Precisão).",
                "Metáfora 4: 'Permeabilidade Seletiva' (Biologia Celular).",
                "Dimensões: Cognitiva, Práxis, Territorial.",
                "Inovação: Avaliação Multidimensional (Não linear)."
            ],
            "notes": "Operacionalizamos com o 'Voxel': o aluno não é uma nota linear, é um ponto 3D. E usamos a 'Permeabilidade Seletiva': a rubrica é uma membrana que deixa entrar o saber local (nutrição) mas mantém o rigor científico (proteção).",
            "image_desc": "Representação 3D de um voxel (cubo) ao lado de uma membrana celular filtrando elementos."
        },
        {
            "layout": 1,
            "title": "O Papel da IA: Dentro do Ecossistema",
            "content": [
                "Metáfora 5: 'O Ecossistema Avaliativo'.",
                "Interdependência: IA, Professores, Alunos, Território.",
                "IA como Tecnologia Assistiva (Visão de Nicolelis).",
                "Papel: Processar dados para liberar o humano."
            ],
            "notes": "Não é software isolado, é um 'Ecossistema'. Os elementos são interdependentes. A IA é tecnologia assistiva: processa a volumetria de dados para liberar o docente para a mediação humana. A IA é o solo fértil, o professor é o jardineiro.",
            "image_desc": "Diagrama de um ecossistema natural interligado com ícones de dados."
        },
        {
            "layout": 1,
            "title": "Funcionalidade: Do Julgamento à Regulação",
            "content": [
                "Mudança: De punição para regulação.",
                "Metacognição: Pensar sobre o próprio pensar.",
                "DUA: Redução de barreiras e ruídos.",
                "Transparência: Critérios explícitos."
            ],
            "notes": "A rubrica funciona como regulação, não punição. Ela reduz o ruído na comunicação. O aluno desenvolve metacognição: entende por que errou. É aprender a aprender, iluminado pela teoria e enraizado no seu chão.",
            "image_desc": "Comparativo: Balança (Justiça Cega) vs. Espelho/Bússola (Orientação)."
        },
        {
            "layout": 1,
            "title": "Metodologia: Design-Based Research",
            "content": [
                "Método: Design-Based Research (DBR).",
                "Fontes: História da EPT (COAGRI/IFTM) e Microdados (INEP).",
                "Validação: Matriz de Rubrica para Avaliar Rubricas (Mullinix).",
                "Tecnologia: Python, Streamlit, LLMs (Grounding)."
            ],
            "notes": "Usamos a DBR: intervir para compreender. Fizemos uma 'Arqueologia Educativa' em Sobradinho. Validamos com uma 'Rubrica para Avaliar Rubricas'. Tecnologicamente, usamos Grounding para evitar alucinações da IA.",
            "image_desc": "Linha do tempo ou fluxograma: Pesquisa -> Modelagem -> Protótipo -> Validação."
        },
        {
            "layout": 1,
            "title": "Resultados: Impacto Social",
            "content": [
                "Transparência Pedagógica.",
                "Combate à Evasão: Monitoramento Preditivo.",
                "Justiça Curricular: No Triângulo Mineiro.",
                "Conceito: Hospitalidade Técnica."
            ],
            "notes": "O resultado vai além da nota. Operamos com 'Hospitalidade Técnica'. A avaliação acolhe as rugosidades do território. Resgatamos a memória das camisetas amarelas da EMS e elevamos à critério pedagógico.",
            "image_desc": "Foto histórica da Fazenda Sobradinho contrastando com dados educacionais futuristas."
        },
        {
            "layout": 1,
            "title": "Futuro: Ecossistema SINAPSE Total",
            "content": [
                "Status: Protótipo Funcional Validado.",
                "Próximos Passos: Mestrado (ProfEPT/MPeT).",
                "Expansão: Ecossistema SINAPSE Total.",
                "Novos Agentes: SOPHIA, YA-YA, MANNA, ROTOR."
            ],
            "notes": "Esta é uma semente. O futuro projeta o 'Ecossistema SINAPSE Total' no mestrado. Novos agentes como SOPHIA e YA-YA atuarão como personas pedagógicas. A IA serve à emancipação humana.",
            "image_desc": "Ícones representando os futuros módulos orbitando o núcleo SINAPSE."
        },
        {
            "layout": 1,
            "title": "Agradecimentos e Referências",
            "content": [
                "Referências Chave: Saviani, Santos, Fernandes, Nicolelis.",
                "Agradecimentos: Orientadora, Família, IFTM.",
                "Acesso: QR Code (GitHub/Streamlit).",
                "Contato: neirivon.cardoso@iftm.edu.br"
            ],
            "notes": "Agradeço à orientadora e colegas. Deixo o link do protótipo à disposição. Estou pronto para as perguntas da banca. Muito obrigado.",
            "image_desc": "Lista limpa de referências e QR Code para o repositório."
        }
    ]
    return slides_data

def create_placeholder_slide(prs, image_desc, left, top, width, height):
    """Cria um retângulo cinza indicando onde a imagem deve entrar."""
    shape = prs.slides[-1].shapes.add_shape(
        1, left, top, width, height  # 1 = msoShapeRectangle
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(200, 200, 200)
    shape.line.color.rgb = RGBColor(100, 100, 100)
    tf = shape.text_frame
    tf.text = f"[INSERIR IMAGEM AQUI]\n{image_desc}"
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)

def main():
    prs = Presentation()
    slides_data = create_slide_content()
    
    # Configurações de tamanho para o placeholder de imagem
    img_left = Inches(5.5)
    img_top = Inches(2.0)
    img_width = Inches(4.0)
    img_height = Inches(3.0)

    for slide_info in slides_data:
        layout = prs.slide_layouts[slide_info["layout"]]
        slide = prs.slides.add_slide(layout)
        
        # Definir Título
        if slide.shapes.title:
            slide.shapes.title.text = slide_info["title"]
        
        # Definir Conteúdo (Bullet Points)
        if slide_info["layout"] == 1 and len(slide.placeholders) > 1:
            body_shape = slide.placeholders[1]
            tf = body_shape.text_frame
            tf.clear()  # Limpar formatação padrão
            
            for point in slide_info["content"]:
                p = tf.add_paragraph()
                p.text = point
                p.level = 0
                p.font.size = Pt(18)
                p.space_after = Pt(10)
        
        # Adicionar Placeholder de Imagem
        create_placeholder_slide(prs, slide_info["image_desc"], img_left, img_top, img_width, img_height)
        
        # Adicionar Notas do Orador
        notes_slide = slide.notes_slide
        text_frame = notes_slide.notes_text_frame
        text_frame.text = slide_info["notes"]

    # Salvar arquivo
    output_filename = "Defesa_TCC_SINAPSE.pptx"
    prs.save(output_filename)
    print(f"Apresentação gerada com sucesso: {output_filename}")

if __name__ == "__main__":
    main()
