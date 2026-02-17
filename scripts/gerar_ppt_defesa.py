"""
================================================================================
ARQUIVO: /home/neirivon/SINAPSE2.0/sinapsebr_rubrica/scripts/gerar_ppt_defesa.py
PROJETO: Rubrica SINAPSE-BR IA
AUTOR: Neirivon Elias Cardoso
ORIENTADORA: Prof. Dra. Thays Martins Vital da Silva
INSTITUIÇÃO: IFTM - Pós-Graduação em Docência para EPT
DATA: 2026-02-16
DESCRIÇÃO: Script para geração automática de esqueleto de slides (.pptx)
            com tópicos e Speaker Notes baseados no roteiro de defesa.
================================================================================
"""

from pptx import Presentation
from pptx.util import Inches, Pt

def create_sinapse_presentation():
    prs = Presentation()

    # Estrutura de dados baseada no Roteiro de Defesa e TCC
    slides_content = [
        {
            "title": "Rubrica SINAPSE-BR IA",
            "bullets": [
                "Sistema Integrado Neuropsicopedagógico de Avaliação",
                "Superando a Dualidade Estrutural na EPT",
                "Integração: Neurociência, Politecnia e Geofilosofia",
                "Autor: Neirivon Elias Cardoso | Orientadora: Thays Martins"
            ],
            "notes": "Apresento a Rubrica SINAPSE-BR IA. O trabalho propõe superar a 'fratura educacional' histórica que divide o pensar do fazer na EPT."
        },
        {
            "title": "O Problema: A Curvatura da Vara",
            "bullets": [
                "Dualidade Estrutural: Pensar vs. Fazer",
                "Avaliação Classificatória vs. Formativa",
                "Metáfora: 'A Curvatura da Vara' (Saviani)",
                "Objetivo: Equilíbrio entre Técnica e Formação Humana"
            ],
            "notes": "O modelo tradicional foca em classificar e dar notas. Para corrigir isso, usamos a metáfora de Saviani: curvar a vara para o lado humano para que ela fique reta."
        },
        {
            "title": "Fundamentação: O Chão e a Luz",
            "bullets": [
                "Geofilosofia (Fernandes, 2023)",
                "Metáfora: 'O Chão Geográfico'",
                "Metáfora: 'A Teoria como Luz' (Freire)",
                "Contexto: Sobradinho e o Triângulo Mineiro"
            ],
            "notes": "A avaliação está enraizada no território. O 'Chão Geográfico' é a realidade social de Sobradinho. A teoria é a chama que ilumina esse chão."
        },
        {
            "title": "A Conexão: Metáfora da Sinapse",
            "bullets": [
                "Nome do Projeto: SINAPSE-BR IA",
                "Metáfora: 'A Sinapse Pedagógica'",
                "Conexão: Ensino + Aprendizagem + Tecnologia",
                "Transformação: Dado -> Informação -> Formação Humana"
            ],
            "notes": "A Sinapse representa o encontro pedagógico. É o momento em que a avaliação vira um feixe de luz que conecta professor, aluno e tecnologia."
        },
        {
            "title": "O Produto: O Voxel e a Permeabilidade",
            "bullets": [
                "Metáfora: 'O Voxel' (Marcador 3D de Precisão)",
                "Metáfora: 'Permeabilidade Seletiva' (Biologia Celular)",
                "Dimensões: Cognitiva, Práxis, Territorial",
                "Avaliação Multidimensional"
            ],
            "notes": "O Voxel posiciona o estudante em um espaço 3D de desenvolvimento. A 'Permeabilidade Seletiva' permite que saberes locais entrem no currículo com rigor científico."
        },
        {
            "title": "IA no Ecossistema Avaliativo",
            "bullets": [
                "Metáfora: 'O Ecossistema Avaliativo'",
                "IA como Tecnologia Assistiva (Nicolelis)",
                "Interdependência: IA, Professores e Alunos",
                "Foco: Mediação Humana e Escuta Pedagógica"
            ],
            "notes": "Não é um software isolado, mas um ecossistema. A IA processa dados complexos para que o docente foque no que é insubstituível: a mediação humana."
        },
        {
            "title": "Metodologia e Validação",
            "bullets": [
                "Design-Based Research (DBR)",
                "Arqueologia Educativa e Documental",
                "Tecnologias: Python, Streamlit e LLMs",
                "Validação: Matriz de Mullinix"
            ],
            "notes": "Analisamos documentos históricos da EMS para validar a dimensão territorial. Tecnologicamente, usamos Python e Streamlit para o protótipo funcional."
        },
        {
            "title": "Resultados: Hospitalidade Técnica",
            "bullets": [
                "Transparência e Combate à Evasão",
                "Justiça Curricular e Territorial",
                "Hospitalidade Técnica (Fernandes, 2023)",
                "Registro Burocrático -> Emancipação Humana"
            ],
            "notes": "A avaliação deixa de ser punitiva para ser acolhedora. Reconhecemos as 'rugosidades do território' que impactam o aprendizado."
        },
        {
            "title": "Futuro: Ecossistema SINAPSE Total",
            "bullets": [
                "Próximos Passos: Mestrado (ProfEPT)",
                "Expansão para Ecossistema Total",
                "Novos Agentes: SOPHIA, YA-YA, MANNA, ROTOR",
                "IA a Serviço da Emancipação Humana"
            ],
            "notes": "Esta é uma semente. Projetamos agentes de IA com personas pedagógicas (SOPHIA, YA-YA) para que a IA sirva humildemente à liberdade humana."
        },
        {
            "title": "Agradecimentos e Referências",
            "bullets": [
                "Principais Referências: Saviani, Santos, Fernandes",
                "Agradecimento especial à Dra. Thays Martins",
                "QR Code para Repositório (GitHub)",
                "Dúvidas e Arguição"
            ],
            "notes": "Agradeço à banca e deixo o repositório à disposição. Estou pronto para as perguntas."
        }
    ]

    for data in slides_content:
        slide_layout = prs.slide_layouts[1] # Layout Título e Conteúdo
        slide = prs.slides.add_slide(slide_layout)
        
        # Título do Slide
        title = slide.shapes.title
        title.text = data["title"]
        
        # Corpo do Slide (Bullets)
        tf = slide.placeholders[1].text_frame
        tf.text = data["bullets"][0]
        for bullet in data["bullets"][1:]:
            p = tf.add_paragraph()
            p.text = bullet
            p.level = 0
            
        # Notas do Orador
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = data["notes"]

    # Salva o arquivo
    filename = "Apresentacao_SINAPSE_BR_IA.pptx"
    prs.save(filename)
    return filename

if __name__ == "__main__":
    try:
        file = create_sinapse_presentation()
        print(f"Sucesso! Apresentação '{file}' gerada com tópicos e Speaker Notes.")
    except Exception as e:
        print(f"Erro ao gerar apresentação: {e}")
